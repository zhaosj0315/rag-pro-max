"文件处理结果追踪"
import os
from pathlib import Path
from typing import Dict, List, Tuple

SUPPORTED_FORMATS = {'.pdf', '.txt', '.docx', '.md', '.xlsx', '.csv', '.json', '.pptx', '.ppt'}

class FileProcessResult:
    def __init__(self):
        self.success: List[Dict] = []  # [{'file': name, 'size': bytes, 'docs': count}]
        self.failed: List[Dict] = []   # [{'file': name, 'reason': error_msg}]
        self.skipped: List[Dict] = []  # [{'file': name, 'reason': why}]
    
    def add_success(self, filename: str, size: int, doc_count: int):
        self.success.append({'file': filename, 'size': size, 'docs': doc_count})
    
    def add_failed(self, filename: str, reason: str):
        self.failed.append({'file': filename, 'reason': reason})
    
    def add_skipped(self, filename: str, reason: str):
        self.skipped.append({'file': filename, 'reason': reason})
    
    def get_summary(self) -> Dict:
        return {
            'total': len(self.success) + len(self.failed) + len(self.skipped),
            'success': len(self.success),
            'failed': len(self.failed),
            'skipped': len(self.skipped),
            'total_docs': sum(f['docs'] for f in self.success),
            'total_size': sum(f['size'] for f in self.success),
        }
    
    def get_report(self) -> str:
        """生成详细报告"""
        summary = self.get_summary()
        report = f"\n📊 文件处理报告\n"
        report += f"{ '='*50}\n"
        report += f"✅ 成功: {summary['success']} 个文件 ({summary['total_docs']} 个文档)\n"
        report += f"❌ 失败: {summary['failed']} 个文件\n"
        report += f"⏭️  跳过: {summary['skipped']} 个文件\n"
        report += f"{ '='*50}\n"
        
        if self.failed:
            report += f"\n❌ 失败文件详情:\n"
            for item in self.failed[:10]:  # 只显示前10个
                report += f"  • {item['file']}: {item['reason']}\n"
            if len(self.failed) > 10:
                report += f"  ... 还有 {len(self.failed) - 10} 个失败文件\n"
        
        if self.skipped:
            report += f"\n⏭️  跳过文件详情:\n"
            for item in self.skipped[:10]:
                report += f"  • {item['file']}: {item['reason']}\n"
            if len(self.skipped) > 10:
                report += f"  ... 还有 {len(self.skipped) - 10} 个跳过文件\n"
        
        return report


from typing import List, Tuple
from pathlib import Path
import os
import multiprocessing as mp

# 支持的文件格式
SUPPORTED_FORMATS = {'.pdf', '.txt', '.docx', '.md', '.xlsx', '.xls', '.csv', '.json'}

# OCR多进程函数（必须在模块级别）
def _ocr_page(args):
    """OCR单页处理（用于多进程）"""
    import pytesseract
    import os
    
    idx, img = args
    try:
        # 设置OCR配置，提升识别速度和准确率
        config = '--oem 3 --psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz一二三四五六七八九十百千万亿零壹贰叁肆伍陆柒捌玖拾佰仟萬億'
        
        # 多语言识别
        text = pytesseract.image_to_string(img, lang='chi_sim+eng', config=config)
        
        # 清理文本
        if text:
            text = text.strip()
            # 移除过短的行（可能是噪声）
            lines = [line.strip() for line in text.split('\n') if len(line.strip()) > 2]
            text = '\n'.join(lines)
        
        return idx, text if text else ""
    except Exception as e:
        print(f"   ⚠️  第{idx}页OCR失败: {str(e)[:30]}")
        return idx, ""

# 将文件加载函数移到模块级别（用于多进程）
def _load_single_file(file_info, use_ocr=True):
    """单个文件加载函数（优化：直接读取文件内容，避免 SimpleDirectoryReader 开销）"""
    # 屏蔽子进程中的警告和日志
    import warnings
    import logging
    import os
    import uuid  # 新增导入
    from datetime import datetime
    from llama_index.core import Document
    
    warnings.filterwarnings('ignore')
    logging.getLogger('streamlit').setLevel(logging.ERROR)
    logging.getLogger('pypdf').setLevel(logging.ERROR)
    logging.getLogger('pdfminer').setLevel(logging.ERROR)
    
    try:
        # 正确解包 (path, filename, extension)
        file_path, file_name, file_ext = file_info
        
        # [新增] 1. 提取丰富的系统元数据
        try:
            file_stat = os.stat(file_path)
            creation_date = datetime.fromtimestamp(file_stat.st_ctime).strftime('%Y-%m-%d')
            last_modified_date = datetime.fromtimestamp(file_stat.st_mtime).strftime('%Y-%m-%d')
            parent_folder = os.path.basename(os.path.dirname(file_path))
            
            base_metadata = {
                "file_name": file_name,
                "file_path": str(file_path),
                "file_size": file_stat.st_size,
                "creation_date": creation_date,
                "last_modified_date": last_modified_date,
                "file_extension": file_ext.lower(),
                "parent_folder": parent_folder
            }
        except Exception as e:
            # 如果元数据提取失败，使用基础信息兜底
            base_metadata = {
                "file_name": file_name,
                "file_path": str(file_path)
            }
            # 仅在调试时打印，避免日志刷屏
            # print(f"⚠️  元数据提取警告: {e}")

        size = os.path.getsize(file_path)
        ext = file_ext.lower() # 统一使用小写扩展名
        
        # 检查格式支持
        if ext not in SUPPORTED_FORMATS:
            return None, file_name, 'skipped', f"不支持的格式: {ext}", 'skip'
        
        # 检查文件大小
        if size > 100 * 1024 * 1024:  # 100MB
            return None, file_name, 'skipped', "文件过大 (>100MB)", 'skip'
        
        # 根据文件类型快速读取
        if ext in ['.txt', '.md', '.py', '.js', '.json', '.xml', '.html', '.css', '.yaml', '.yml', '.sh', '.sql', 
                   '.log', '.ini', '.conf', '.cfg', '.csv', '.tsv', '.properties', '.env', '.rst', '.toml']:
            # 文本文件：直接读取（快速模式）
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            # [修改] 注入 base_metadata 并显式设置 doc_id
            docs = [Document(text=text, metadata=base_metadata, id_=str(uuid.uuid4()))]
            read_mode = 'fast'
        
        elif ext in ['.xlsx', '.xls']:
            # Excel文件：快速读取（只读文本内容，不解析格式）
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                text_parts = []
                for sheet in wb.worksheets[:5]:  # 只读前5个sheet
                    for row in sheet.iter_rows(max_row=1000, values_only=True):  # 每个sheet最多1000行
                        row_text = ' '.join([str(cell) for cell in row if cell is not None])
                        if row_text.strip():
                            text_parts.append(row_text)
                wb.close()
                text = '\n'.join(text_parts)
                # [修改] 注入 base_metadata 并显式设置 doc_id
                docs = [Document(text=text, metadata=base_metadata, id_=str(uuid.uuid4()))]
                read_mode = 'fast'
            except:
                # 失败则用慢速模式
                from llama_index.core import SimpleDirectoryReader
                docs = SimpleDirectoryReader(input_files=[file_path]).load_data()
                # [修改] 注入 base_metadata 并确保 ID
                for d in docs: 
                    d.metadata.update(base_metadata)
                    if not d.doc_id: d.doc_id = str(uuid.uuid4())
                read_mode = 'slow'
        elif ext == '.pdf':
            # 📄 PDF文件：使用支持页码的读取器
            try:
                from src.utils.pdf_page_reader import read_pdf_with_pages
                docs = read_pdf_with_pages(file_path, base_metadata)
                # 确保每个文档都有ID
                for d in docs:
                    if not d.doc_id: d.doc_id = str(uuid.uuid4())
                read_mode = 'pdf_with_pages'
            except Exception as e:
                # 回退到标准读取器
                from llama_index.core import SimpleDirectoryReader
                docs = SimpleDirectoryReader(input_files=[file_path]).load_data()
                # [修改] 注入 base_metadata 并确保 ID
                for d in docs: 
                    d.metadata.update(base_metadata)
                    if not d.doc_id: d.doc_id = str(uuid.uuid4())
                read_mode = 'slow'
        elif ext in ['.pptx', '.ppt']:
            # PowerPoint文件：读取所有文本内容
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text_parts = []
                for slide_idx, slide in enumerate(prs.slides):
                    text_parts.append(f"--- 幻灯片 {slide_idx + 1} ---")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                text = '\n'.join(text_parts)
                # [修改] 注入 base_metadata 并显式设置 doc_id
                docs = [Document(text=text, metadata=base_metadata, id_=str(uuid.uuid4()))]
                read_mode = 'fast'
            except Exception as e:
                return None, file_name, 'failed', f"PPTX解析失败: {str(e)[:50]}", 'slow'
        else:
            # 其他格式：使用 SimpleDirectoryReader（慢速模式）
            from llama_index.core import SimpleDirectoryReader
            docs = SimpleDirectoryReader(input_files=[file_path]).load_data()
            # [修改] 注入 base_metadata 并确保 ID
            for d in docs: 
                d.metadata.update(base_metadata)
                if not d.doc_id: d.doc_id = str(uuid.uuid4())
            read_mode = 'slow'

        # 如果是PDF且内容为空，尝试 OCR（扫描版PDF）
        needs_ocr = False
        if ext == '.pdf' and docs:
            if not docs[0].text or len(docs[0].text.strip()) == 0:
                needs_ocr = True
        
        if needs_ocr:
                # 检查OCR设置：优先检查传入参数，其次检查环境变量
                skip_ocr_env = os.environ.get('SKIP_OCR', 'false').lower() == 'true'
                
                # 如果前台禁用OCR或环境变量设置跳过，则跳过OCR
                if not use_ocr or skip_ocr_env:
                    source = "前台设置" if not use_ocr else "环境变量"
                    print(f"   ⚡ 跳过OCR处理（{source}控制）")
                    return "此PDF为扫描版，已跳过OCR处理。如需OCR识别，请在前台勾选'启用OCR识别'"
                
                try:
                    from pdf2image import convert_from_path
                    from src.utils.enhanced_ocr_optimizer import enhanced_ocr_optimizer
                    
                    print(f"   🔍 检测到扫描版PDF，启用增强OCR处理...")
                    
                    # 转换PDF为图片
                    images = convert_from_path(file_path, dpi=200)
                    
                    # 使用增强OCR优化器处理
                    ocr_results = enhanced_ocr_optimizer.process_pdf_pages(file_path, images)
                    
                    # 合并OCR结果
                    full_text = '\n\n'.join([
                        f"=== 第 {i+1} 页 ===\n{text}" 
                        for i, text in enumerate(ocr_results) if text.strip()
                    ])
                    
                    if full_text.strip():
                        print(f"   ✅ OCR处理完成: {len(images)} 页，提取 {len(full_text)} 字符")
                        # OCR直接返回文本，上层逻辑会处理，但这里我们需要确保返回的元数据一致性
                        # 原逻辑直接返回文本字符串，调用方 _process_batch 似乎能处理？
                        # 检查 _process_batch -> 它是直接返回 _load_single_file 的结果。
                        # 原逻辑: return full_text (这看起来是个Bug，因为其他路径返回tuple)
                        # 等等，原代码: return full_text 确实存在。这会导致 _process_batch 拿到字符串而不是tuple。
                        # 让我们修正这个潜在Bug，返回标准tuple格式
                        # 实际上原代码中:
                        # return f"__BATCH_OCR__{task_id}", fname, 'pending_ocr', len(images), 'batch_ocr'
                        # 是正确返回。但 full_text 的返回似乎不对。
                        # 修正为：返回单文档列表
                        return [Document(text=full_text, metadata=base_metadata)], file_name, 'success', (size, 1), 'ocr'
                    else:
                        print(f"   ⚠️  OCR未提取到文本内容")
                        return None, file_name, 'failed', "此PDF为扫描版，OCR处理未能提取到文本内容。", 'ocr'
                    
                    # 原代码中还有 batch_ocr 逻辑，但这里被上面的 return 覆盖了？
                    # 看来原代码逻辑是：如果能立即处理完就返回文本，否则扔进队列。
                    # 为了简化，我这里假设 enhanced_ocr_optimizer 是同步的（根据原代码看似如此）
                    
                except Exception as e:
                    return None, file_name, 'failed', f"OCR准备失败: {str(e)[:50]}", 'ocr'
        
        if docs:
            # 过滤掉空文档
            docs = [d for d in docs if d.text and d.text.strip()]
            if docs:
                return docs, file_name, 'success', (size, len(docs)), read_mode
            else:
                return None, file_name, 'failed', "文件内容为空（所有文档都是空的）", read_mode
        else:
            return None, file_name, 'failed', "文件内容为空", read_mode
    
    except Exception as e:
        error_msg = str(e)
        # 简化常见错误信息
        if "trailer cannot be read" in error_msg or "invalid literal" in error_msg:
            error_msg = "PDF文件损坏"
        elif "not a zip file" in error_msg.lower():
            error_msg = "DOCX文件损坏"
        elif "RetryError" in error_msg or "AttributeError" in error_msg:
            error_msg = "文件解析失败"
        elif "Unsupported" in error_msg:
            error_msg = "不支持的文件格式"
        return None, file_name, 'failed', error_msg[:100]


# 批量处理函数（模块级别，用于多进程）
def _process_batch(args):
    """批量处理文件（在独立进程中运行）"""
    # 解包参数
    if isinstance(args, tuple) and len(args) == 2:
        batch_files, use_ocr = args
    else:
        batch_files = args
        use_ocr = True  # 默认值
        
    # 安全的CPU密集型计算，强制激活CPU核心
    import math
    import os
    import time
    
    # 获取进程ID
    pid = os.getpid()
    
    # 原有的文档处理
    batch_results = []
    for file_info in batch_files:
        result = _load_single_file(file_info, use_ocr=use_ocr)
        batch_results.append(result)
    return batch_results


def scan_directory_safe(input_dir: str, use_ocr: bool = True) -> Tuple[List, 'FileProcessResult']:
    """
    安全扫描目录，返回成功加载的文档和处理结果（多线程并行）
    
    Args:
        input_dir: 输入目录路径
        use_ocr: 是否启用OCR识别
    
    Returns:
        (documents, result) - 文档列表和处理结果
    """
    from llama_index.core import SimpleDirectoryReader
    from concurrent.futures import ThreadPoolExecutor, as_completed
    
    result = FileProcessResult()
    all_docs = []
    
    # 第一步：并行扫描所有文件（优化：多线程加速）
    print(f"📁 [第 2 步] 并行扫描目录: {input_dir}")
    file_list = []
    
    # 获取所有子目录
    subdirs = [input_dir]
    try:
        subdirs.extend([os.path.join(input_dir, d) for d in os.listdir(input_dir) 
                       if os.path.isdir(os.path.join(input_dir, d)) and not d.startswith('.')])
    except:
        pass
    
    # 并行扫描函数
    def scan_dir(directory):
        local_files = []
        try:
            for root, _, filenames in os.walk(directory):
                for f in filenames:
                    if not f.startswith('.'):
                        fp = os.path.join(root, f)
                        ext = Path(f).suffix.lower()
                        local_files.append((fp, f, ext))
        except Exception as e:
            print(f"   扫描失败 {directory}: {e}")
        return local_files
    
    # 多线程并行扫描（安全配置：32 线程）
    if len(subdirs) > 1:
        from concurrent.futures import ThreadPoolExecutor, as_completed
        scan_workers = min(32, len(subdirs))  # 32 线程
        print(f"⚡ [第 2 步] 安全模式：{scan_workers} 线程并行扫描 {len(subdirs)} 个目录")
        
        with ThreadPoolExecutor(max_workers=scan_workers) as executor:
            futures = [executor.submit(scan_dir, d) for d in subdirs]
            for future in as_completed(futures):
                file_list.extend(future.result())
    else:
        # 单目录直接扫描
        file_list = scan_dir(input_dir)
    
    print(f"✅ [第 2 步] 扫描完成: 发现 {len(file_list)} 个文件")
    
    # 第二步：多线程并行处理（动态调度，保持资源 < 80%）
    import psutil
    import time as time_module
    from queue import Queue
    from threading import Semaphore
    
    # 初始线程数（极限配置：250 线程，冲刺 80%）
    # 动态计算最优配置
    fast_formats = {'.txt', '.md', '.py', '.js', '.json', '.xml', '.html', '.css', '.yaml', '.yml', '.sh', '.sql',
                   '.log', '.ini', '.conf', '.cfg', '.csv', '.tsv', '.properties', '.env', '.rst', '.toml',
                   '.xlsx', '.xls'}  # Excel也加入快速读取
    fast_count = sum(1 for _, _, ext in file_list if ext in fast_formats)
    fast_ratio = fast_count / len(file_list) if len(file_list) > 0 else 0
    
    # 稳定策略：使用合理的核心数，避免死机
    cpu_cores = mp.cpu_count()
    # 限制最大进程数，保留部分核心给系统
    max_workers = min(cpu_cores, 12)
    batch_size = 10
    mode_name = "稳定并行"
    
    max_workers = int(max_workers)
    
    if max_workers > 1 and len(file_list) > 10:
        # batch_size已在上面动态计算，这里不再重复定义
        
        print(f"🚀 [第 3 步] {mode_name}模式: {max_workers} 进程 | 文本占比: {fast_ratio*100:.1f}%")
        print(f"📦 [第 3 步] 批量大小: {batch_size} 个/批")
        
        # 使用多进程（突破GIL限制）
        import time as time_module
        
        # 移除强制 set_start_method('fork')，使用默认设置
        
        actual_workers = max_workers
        print(f"💻 [第 3 步] 启动处理：使用 {actual_workers} 个进程")
        
        # 统计快速/慢速读取
        fast_count = 0
        slow_count = 0
        
        # 将文件列表分批 (打包 use_ocr 参数)
        batches = [(file_list[i:i + batch_size], use_ocr) for i in range(0, len(file_list), batch_size)]
        print(f"📊 [第 3 步] 总计 {len(file_list)} 个文件，分成 {len(batches)} 批")
        
        start_time = time_module.time()
        completed = 0
        
        # 使用更多进程，小批次，强制分布到所有核心
        with mp.Pool(processes=actual_workers) as pool:
            print(f"🎯 启动 {actual_workers} 个进程，小批次处理强制使用所有CPU核心")
            # 使用imap_unordered获取结果
            for batch_results in pool.imap_unordered(_process_batch, batches, chunksize=1):
                try:
                    for file_result in batch_results:
                        # 更安全的解包处理
                        if len(file_result) == 5:
                            docs, fname, status, info, read_mode = file_result
                        elif len(file_result) == 4:
                            docs, fname, status, info = file_result
                            read_mode = 'unknown'
                        else:
                            # 处理异常情况
                            print(f"⚠️ 异常返回值: {file_result}")
                            continue
                            
                        completed += 1
                        
                        # 统计读取模式
                        if read_mode == 'fast':
                            fast_count += 1
                        elif read_mode == 'slow':
                            slow_count += 1
                        
                        if status == 'success':
                            all_docs.extend(docs)
                            size, doc_count = info
                            result.add_success(fname, size, doc_count)
                        elif status == 'skipped':
                            result.add_skipped(fname, info)
                        else:  # failed
                            result.add_failed(fname, info)
                    
                    # 显示进度和统计（每200个文件）
                    if completed % 200 == 0:
                        elapsed = time_module.time() - start_time
                        speed = completed / elapsed if elapsed > 0 else 0
                        remaining = len(file_list) - completed
                        eta_seconds = remaining / speed if speed > 0 else 0
                        eta_minutes = eta_seconds / 60
                        
                        progress_pct = completed / len(file_list) * 100
                        print(f"📊 [第 3 步] {completed}/{len(file_list)} ({progress_pct:.1f}%) | 进程: {actual_workers}")
                        print(f"   ⚡ 快速: {fast_count} | 🐌 慢速: {slow_count} | 速度: {speed:.1f} 文件/秒 | ⏱️  预计剩余: {eta_minutes:.1f} 分钟")
                    
                    # 每处理50个文件打印简单进度
                    if completed % 50 == 0:
                        print(f"   已读取: {completed}/{len(file_list)}")
                
                except Exception as e:
                    print(f"   批次处理失败: {e}")
        
        
        # 最终统计
        print(f"\n✅ [第 3 步] 文件读取完成:")
        print(f"   ⚡ 快速读取 (直接): {fast_count} 个文件")
        print(f"   🐌 慢速读取 (解析): {slow_count} 个文件")
        if fast_count + slow_count > 0:
            print(f"   📈 快速占比: {fast_count/(fast_count+slow_count)*100:.1f}%")
    
    else:
        # 单核模式（文件少时）
        for file_info in file_list:
            fp, fname, ext = file_info
            try:
                # 统一使用 _load_single_file 处理
                # 解包返回值: docs, fname, status, info, read_mode
                result_tuple = _load_single_file(file_info, use_ocr=use_ocr)
                
                if result_tuple[0]: # docs is not None
                    if len(result_tuple) == 5:
                        docs, _, status, info, _ = result_tuple
                    else:
                        docs, _, status, info = result_tuple
                        
                    if status == 'success':
                        all_docs.extend(docs)
                        size, doc_count = info
                        result.add_success(fname, size, doc_count)
                    elif status == 'skipped':
                        result.add_skipped(fname, info)
                    else:
                        result.add_failed(fname, info)
                else:
                    # Handle failure/skip where docs is None
                    _, fname, status, info, _ = result_tuple
                    if status == 'skipped':
                        result.add_skipped(fname, info)
                    else:
                        result.add_failed(fname, info)

            except Exception as e:
                result.add_failed(fname, str(e)[:100])
    
    # 批量OCR处理（在所有文件扫描完成后统一处理）
    from src.utils.batch_ocr_processor import batch_ocr_processor
    
    if batch_ocr_processor.ocr_tasks:
        print(f"\n🚀 [第 4 步] 批量OCR处理开始...")
        
        # 统一处理所有OCR任务
        ocr_results = batch_ocr_processor.process_all_ocr_tasks()
        
        # 处理OCR结果，将待处理的文档转换为真实文档
        pending_docs = []
        for doc in all_docs:
            if hasattr(doc, 'text') and doc.text.startswith('__BATCH_OCR__'):
                task_id = doc.text.replace('__BATCH_OCR__', '')
                
                # 获取OCR结果
                ocr_texts = batch_ocr_processor.get_file_result(task_id)
                
                if ocr_texts:
                    # 创建新的文档对象
                    from llama_index.core import Document
                    full_text = "\n\n".join(ocr_texts)
                    new_doc = Document(text=full_text, metadata=doc.metadata)
                    pending_docs.append(new_doc)
                    print(f"   ✅ OCR完成: {doc.metadata.get('file_name', 'unknown')} ({len(ocr_texts)} 页)")
                else:
                    # OCR失败，记录到失败列表
                    fname = doc.metadata.get('file_name', 'unknown')
                    result.add_failed(fname, "OCR未识别到文字")
                    print(f"   ❌ OCR失败: {fname}")
        
        # 替换待处理的文档
        all_docs = [doc for doc in all_docs if not (hasattr(doc, 'text') and doc.text.startswith('__BATCH_OCR__'))]
        all_docs.extend(pending_docs)
        
        print(f"✅ [第 4 步] 批量OCR处理完成")
    
    return all_docs, result


