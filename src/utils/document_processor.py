from src.app_logging.log_manager import LogManager

logger = LogManager()

"""
文档处理模块 - 统一管理文档加载、解析和元数据提取
"""
import os
from datetime import datetime
from llama_index.core.schema import Document
from src.common.utils import sanitize_filename



def get_file_size_str(size_bytes: int) -> str:
    """将字节数转换为可读的文件大小字符串"""
    if size_bytes < 1024 * 1024:
        return f"{size_bytes/1024:.1f} KB"
    else:
        return f"{size_bytes/(1024*1024):.1f} MB"


def get_file_type(filename: str) -> tuple:
    """
    根据文件扩展名返回文件类型和图标
    
    Returns:
        tuple: (类型名称, 图标emoji)
    """
    ext = os.path.splitext(filename)[1].lower()
    
    if ext in ['.pdf']:
        return "PDF", "📄"
    elif ext in ['.docx', '.doc', '.rtf']:
        return "DOC", "📝"
    elif ext in ['.pptx', '.ppt', '.odp']:
        return "PPTX", "🎯"
    elif ext in ['.xls', '.xlsx', '.csv']:
        return "DATA", "📊"
    elif ext in ['.md', '.txt', '.log', '.json', '.xml']:
        return "TEXT", "📜"
    elif ext in ['.jpg', '.png', '.jpeg', '.gif']:
        return "IMG", "🖼️"
    elif ext in ['.zip']:
        return "ZIP", "📦"
    else:
        return "OTHER", "💡"


def load_pptx_file(file_path: str) -> list:
    """
    加载 PPTX 文件
    
    Args:
        file_path: PPTX 文件路径
    
    Returns:
        list: Document 对象列表
    """
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_content = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = f"--- 幻灯片 {slide_idx + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            text_content.append(slide_text)
        
        full_text = "\n".join(text_content)
        return [Document(text=full_text, metadata={"file_path": file_path, "file_type": "pptx"})]
    except Exception as e:
        logger.error(e)
        return []


def get_file_info(file_path: str, metadata_mgr=None) -> dict:
    """
    获取文件的基本信息和元数据
    
    Args:
        file_path: 文件路径
        metadata_mgr: 元数据管理器（可选）
    
    Returns:
        dict: 文件信息字典
    """
    try:
        size_bytes = os.path.getsize(file_path)
        size_str = get_file_size_str(size_bytes)
        file_name = os.path.basename(file_path)
        file_type, file_icon = get_file_type(file_name)
        
        # [新增] 提取更多系统元数据
        file_stat = os.stat(file_path)
        creation_date = datetime.fromtimestamp(file_stat.st_ctime).strftime('%Y-%m-%d')
        last_modified = datetime.fromtimestamp(file_stat.st_mtime).strftime('%Y-%m-%d')
        parent_folder = os.path.basename(os.path.dirname(file_path))
        
        info = {
            "name": file_name,
            "file_path": os.path.abspath(file_path),
            "size": size_str,
            "size_bytes": size_bytes,
            "added_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
            "creation_date": creation_date,
            "last_modified": last_modified,
            "parent_folder": parent_folder,
            "type": file_type,
            "icon": file_icon
        }
        
        # 尝试从 metadata_mgr 获取更多智能信息
        if metadata_mgr:
            meta = metadata_mgr.get_metadata(file_name)
            if meta:
                info.update(meta)
                
        return info
    except Exception as e:
        logger.error(e)
        return {
            "name": os.path.basename(file_path),
            "size": "0 KB",
            "type": "UNKNOWN",
            "icon": "❓",
            "error": str(e)
        }


def get_relevance_label(score: float) -> str:
    """
    根据相关性分数返回标签
    
    Args:
        score: 相关性分数 (0-1)
    
    Returns:
        str: 相关性标签
    """
    if score > 0.8:
        return "🟢 高"
    elif score > 0.7:
        return "🟡 中"
    else:
        return "⚪️ 低"